#!/usr/bin/env python3
"""
python-pptx 演示脚本
展示如何创建、编辑 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

OUTPUT_DIR = "/root/data/disk/workspace"


def create_title_slide():
    """创建标题页"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI 生成报告"
    subtitle.text = "2024年度总结\n由 python-pptx 自动生成"
    
    path = os.path.join(OUTPUT_DIR, "title_slide.pptx")
    prs.save(path)
    print(f"✅ 标题页已保存: {path}")
    return path


def create_content_slide():
    """创建内容页"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "项目成果"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "主要成就："
    
    p = tf.add_paragraph()
    p.text = "✓ 用户增长 150%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ 收入突破 1 亿"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ 团队扩展至 50 人"
    p.level = 1
    
    path = os.path.join(OUTPUT_DIR, "content_slide.pptx")
    prs.save(path)
    print(f"✅ 内容页已保存: {path}")
    return path


def create_table_slide():
    """创建表格页"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
    title_box.text_frame.text = "季度数据"
    
    # 创建表格
    rows, cols = 4, 3
    left, top, width, height = Inches(1.5), Inches(1.5), Inches(6), Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 表头
    table.cell(0, 0).text = "季度"
    table.cell(0, 1).text = "收入(万)"
    table.cell(0, 2).text = "增长率"
    
    # 数据
    data = [
        ("Q1", "120", "15%"),
        ("Q2", "150", "25%"),
        ("Q3", "180", "20%"),
    ]
    
    for i, (quarter, revenue, growth) in enumerate(data, 1):
        table.cell(i, 0).text = quarter
        table.cell(i, 1).text = revenue
        table.cell(i, 2).text = growth
    
    path = os.path.join(OUTPUT_DIR, "table_slide.pptx")
    prs.save(path)
    print(f"✅ 表格页已保存: {path}")
    return path


def create_chart_slide():
    """创建图表页"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
    title_box.text_frame.text = "月度销售趋势"
    
    # 创建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = ['1月', '2月', '3月', '4月', '5月', '6月']
    chart_data.add_series('2024', (30, 45, 52, 48, 65, 70))
    chart_data.add_series('2025', (40, 55, 68, 72, 85, 90))
    
    # 添加图表
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    path = os.path.join(OUTPUT_DIR, "chart_slide.pptx")
    prs.save(path)
    print(f"✅ 图表页已保存: {path}")
    return path


def create_multi_slide():
    """创建完整演示文稿（多页）"""
    prs = Presentation()
    
    # 第1页：标题
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "2024 年度总结"
    slide.placeholders[1].text = "技术部"
    
    # 第2页：目录
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录"
    slide.placeholders[1].text = "1. 业绩回顾\n2. 产品进展\n3. 团队建设\n4. 未来规划"
    
    # 第3页：内容
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "业绩回顾"
    slide.placeholders[1].text = "全年收入增长 45%，用户数突破 100 万"
    
    # 第4页：空白页（可自定义）
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_box.text_frame.text = "谢谢！"
    title_box.text_frame.paragraphs[0].font.size = Pt(44)
    title_box.text_frame.paragraphs[0].font.bold = True
    
    path = os.path.join(OUTPUT_DIR, "full_presentation.pptx")
    prs.save(path)
    print(f"✅ 完整演示文稿已保存: {path}")
    return path


if __name__ == "__main__":
    print("🎨 python-pptx 演示开始...\n")
    
    create_title_slide()
    create_content_slide()
    create_table_slide()
    create_chart_slide()
    create_multi_slide()
    
    print("\n✨ 所有演示文件已生成！")
