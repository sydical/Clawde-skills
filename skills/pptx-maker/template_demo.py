#!/usr/bin/env python3
"""
使用模板创建 PPT
演示如何基于现有模板生成定制化 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os

OUTPUT_DIR = "/root/data/disk/workspace"


def create_template_demo():
    """
    演示如何使用模板
    实际使用时需要先准备一个 template.pptx 文件
    """
    template_path = os.path.join(OUTPUT_DIR, "template.pptx")
    
    # 检查模板是否存在
    if not os.path.exists(template_path):
        print("⚠️ 模板文件不存在，将创建一个空白演示文稿作为模板")
        # 创建一个空白模板
        prs = Presentation()
        prs.save(template_path)
        print(f"✅ 空白模板已创建: {template_path}")
    
    # 使用模板创建新 PPT
    prs = Presentation(template_path)
    
    # 添加新幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # 修改标题
    slide.shapes.title.text = "基于模板的内容"
    
    # 修改正文
    slide.placeholders[1].text = "这是使用模板生成的内容"
    
    # 保存
    output_path = os.path.join(OUTPUT_DIR, "from_template.pptx")
    prs.save(output_path)
    print(f"✅ 基于模板的 PPT 已保存: {output_path}")
    return output_path


def create_business_template():
    """
    创建一个商务风格模板并使用
    """
    # 创建演示文稿
    prs = Presentation()
    
    # 设置封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "商务报告"
    slide.placeholders[1].text = "2024年度总结"
    
    # 添加一页内容
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录"
    slide.placeholders[1].text = "1. 公司概况\n2. 业务进展\n3. 财务数据\n4. 发展规划"
    
    # 保存模板
    template_path = os.path.join(OUTPUT_DIR, "business_template.pptx")
    prs.save(template_path)
    print(f"✅ 商务模板已创建: {template_path}")
    
    # 使用模板
    prs2 = Presentation(template_path)
    
    # 添加新内容
    slide = prs2.slides.add_slide(prs2.slide_layouts[1])
    slide.shapes.title.text = "公司概况"
    slide.placeholders[1].text = "公司成立于2020年，专注于AI技术创新"
    
    output_path = os.path.join(OUTPUT_DIR, "business_report.pptx")
    prs2.save(output_path)
    print(f"✅ 商务报告已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    print("📋 模板演示开始...\n")
    
    create_template_demo()
    print()
    create_business_template()
    
    print("\n✨ 模板演示完成！")
